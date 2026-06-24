# -*- coding: utf-8 -*-
"""Build the CEBA Procurement Academy deck:
"Understanding Vietnam's DPPA Mechanisms and Pricing Considerations"

45-minute Day 2 technical session (virtual / grid CfD DPPA only).
Run:  python outputs/ceba_training/build_deck.py
Output: outputs/ceba_training/CEBA_DPPA_Mechanisms_Pricing.pptx

All figures are hard-coded constants below with their repo source so a
refresh is a one-place edit.
"""

import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ---------------------------------------------------------------------------
# Style (matched to Factory_A_Solar_BESS_Case_Study.pptx)
# ---------------------------------------------------------------------------
TEAL = RGBColor(0x12, 0x7A, 0x7A)
NAVY = RGBColor(0x0F, 0x29, 0x3D)
GREEN = RGBColor(0x2E, 0x9E, 0x5B)
AMBER = RGBColor(0xE0, 0x9A, 0x2B)
INK = RGBColor(0x22, 0x2B, 0x33)
GREY = RGBColor(0x5A, 0x66, 0x70)
PALE = RGBColor(0xEE, 0xF4, 0xF4)      # pale teal fill
PALE2 = RGBColor(0xF5, 0xF6, 0xF7)     # pale grey fill
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RED = RGBColor(0xB3, 0x3A, 0x3A)

TITLE_FONT = "Segoe UI Semibold"
BODY_FONT = "Segoe UI"
MONO_FONT = "Consolas"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ---------------------------------------------------------------------------
# Source figures
# ---------------------------------------------------------------------------
# vietnam_market_context.md - EVN TOU structure & rates (2025, 25,000 VND/USD)
TOU_PERIODS = [
    ("Off-peak (thap diem)", "22:00 - 04:00", "~0.52x"),
    ("Normal (binh thuong)", "04:00-09:30, 11:30-17:00, 20:00-22:00", "1.0x"),
    ("Peak (cao diem)", "09:30-11:30, 17:00-20:00", "~1.78x"),
]
TOU_RATES = [
    (">=110kV", "0.034", "0.065", "0.116"),
    ("22-110kV", "0.037", "0.070", "0.126"),
    ("6-22kV", "0.042", "0.079", "0.141"),
    ("<6kV", "0.049", "0.092", "0.164"),
]
# vietnam_market_context.md - K_pp official 2025
KPP_ROWS = [
    (">=110kV", "1.008525", "2.79%", "Eligible"),
    ("22-110kV", "1.027263", "4.56%", "Eligible"),
    ("6-22kV", "TBD (NLDC/EVN)", "-", "Not yet usable"),
    ("<6kV", "TBD (NLDC/EVN)", "-", "Not yet usable"),
]
# vietnam_market_context.md - FMP/CFMP 2025 means
FMP_2025 = "1,426.6"
CFMP_2025 = "1,464.8"
# NSMO CD7 deck (tmp_dppa_simulation_extract.txt) - Vi du 1 assumptions
VD1 = {
    "k": "1.026", "kpp": "1.008", "cdppa": "360", "pcl": "163.3",
    "p1": "2,204", "pc": "1,300", "fmp": "1,200",
}
# SESSION_NOTES.md 2026-06-11 corrected results
CASE5 = {
    "irr": "16.9%", "pirr": "13.5%", "npv": "$1.52M", "dscr": "1.14x",
    "payback": "9.1 yr", "b1": "-8.7%", "b10": "-8.9%", "blife": "-9.3%",
}
CASE6 = {
    "irr": "26.9%", "pirr": "18.2%", "npv": "$2.54M", "dscr": "1.50x",
    "payback": "4.7 yr", "b1": "-14.4%", "b10": "-14.4%", "blife": "-14.4%",
}

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]

_slide_no = 0


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------
def _set_text(tf, lines, font=BODY_FONT, size=14, color=INK, bold=False,
              align=PP_ALIGN.LEFT, line_spacing=1.0, space_after=4):
    """lines: list of (level, text, overrides-dict) or plain strings."""
    tf.word_wrap = True
    first = True
    for item in lines:
        ov = {}
        if isinstance(item, tuple):
            if isinstance(item[0], int):
                level, text = item[0], item[1]
                if len(item) > 2:
                    ov = item[2]
            else:
                level, text = 0, item[0]
                if len(item) > 1:
                    ov = item[1]
        else:
            level, text = 0, item
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = level
        p.alignment = ov.get("align", align)
        p.line_spacing = line_spacing
        p.space_after = Pt(space_after)
        run = p.add_run()
        run.text = text
        f = run.font
        f.name = ov.get("font", font)
        f.size = Pt(ov.get("size", size))
        f.bold = ov.get("bold", bold)
        f.color.rgb = ov.get("color", color)


def textbox(slide, x, y, w, h, lines, **kw):
    box = slide.shapes.add_textbox(x, y, w, h)
    _set_text(box.text_frame, lines, **kw)
    return box


def content_slide(kicker, title, accent=TEAL, footer=True):
    global _slide_no
    _slide_no += 1
    slide = prs.slides.add_slide(BLANK)
    textbox(slide, Inches(0.55), Inches(0.30), Inches(12.2), Inches(0.32),
            [kicker], font=TITLE_FONT, size=12, color=accent, bold=True)
    textbox(slide, Inches(0.55), Inches(0.62), Inches(12.2), Inches(0.75),
            [title], font=TITLE_FONT, size=26, color=NAVY, bold=True)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.58),
                                 Inches(1.42), Inches(1.6), Pt(3))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    if footer:
        textbox(slide, Inches(0.55), Inches(7.08), Inches(9.5), Inches(0.3),
                ["Understanding Vietnam's DPPA Mechanisms & Pricing  -  CEBA Procurement Academy 2026"],
                size=9, color=GREY)
        textbox(slide, Inches(12.45), Inches(7.08), Inches(0.6), Inches(0.3),
                [str(_slide_no)], size=9, color=GREY, align=PP_ALIGN.RIGHT)
    return slide


def add_table(slide, x, y, w, headers, rows, col_widths=None, font_size=12,
              header_fill=NAVY, row_h=Inches(0.34), header_size=None,
              highlight_rows=(), highlight_fill=PALE):
    shape = slide.shapes.add_table(len(rows) + 1, len(headers), x, y, w,
                                   row_h * (len(rows) + 1))
    tbl = shape.table
    if col_widths:
        total = sum(col_widths)
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = Emu(int(w * cw / total))
    for j, htext in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_fill
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_top = cell.margin_bottom = Pt(2)
        _set_text(cell.text_frame, [htext], font=TITLE_FONT,
                  size=header_size or font_size, color=WHITE, bold=True,
                  space_after=0)
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = tbl.cell(i + 1, j)
            cell.fill.solid()
            if i in highlight_rows:
                cell.fill.fore_color.rgb = highlight_fill
            else:
                cell.fill.fore_color.rgb = WHITE if i % 2 == 0 else PALE2
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_top = cell.margin_bottom = Pt(2)
            if isinstance(val, tuple):
                text, ov = val
            else:
                text, ov = val, {}
            _set_text(cell.text_frame, [(0, text, ov)], font=BODY_FONT,
                      size=font_size, color=INK, space_after=0)
    return tbl


def formula_box(slide, x, y, w, h, lines, size=13):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    box.fill.solid()
    box.fill.fore_color.rgb = PALE
    box.line.color.rgb = TEAL
    box.line.width = Pt(1)
    tf = box.text_frame
    tf.margin_left = tf.margin_right = Inches(0.18)
    tf.margin_top = tf.margin_bottom = Inches(0.08)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    _set_text(tf, lines, font=MONO_FONT, size=size, color=NAVY,
              space_after=2)
    return box


def callout(slide, x, y, w, h, lines, fill=NAVY, color=WHITE, size=14):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    box.line.fill.background()
    tf = box.text_frame
    tf.margin_left = tf.margin_right = Inches(0.2)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    _set_text(tf, lines, font=TITLE_FONT, size=size, color=color, bold=True,
              space_after=2)
    return box


def flow_box(slide, x, y, w, h, title_text, sub, fill):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    box.line.fill.background()
    tf = box.text_frame
    tf.margin_left = tf.margin_right = Inches(0.08)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    lines = [(0, title_text, {"font": TITLE_FONT, "size": 14, "bold": True,
                              "color": WHITE, "align": PP_ALIGN.CENTER})]
    for s in sub:
        lines.append((0, s, {"size": 10.5, "color": WHITE,
                             "align": PP_ALIGN.CENTER}))
    _set_text(tf, lines, space_after=1)
    return box


def arrow_label(slide, x, y, w, text, color=GREY, size=10.5, align=PP_ALIGN.CENTER):
    textbox(slide, x, y, w, Inches(0.5), [text], size=size, color=color,
            align=align)


def connector(slide, x1, y1, x2, y2, color=GREY, weight=2.0, dashed=False):
    conn = slide.shapes.add_connector(2, x1, y1, x2, y2)  # 2 = straight
    conn.line.color.rgb = color
    conn.line.width = Pt(weight)
    if dashed:
        from pptx.oxml.ns import qn
        ln = conn.line._get_or_add_ln()
        dash = ln.makeelement(qn("a:prstDash"), {"val": "dash"})
        ln.append(dash)
    return conn


# ===========================================================================
# SLIDE 1 - Title
# ===========================================================================
_slide_no += 1
s = prs.slides.add_slide(BLANK)
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.18))
band.fill.solid(); band.fill.fore_color.rgb = TEAL; band.line.fill.background()
textbox(s, Inches(0.8), Inches(1.45), Inches(11.7), Inches(0.4),
        ["CEBA CLEAN ENERGY PROCUREMENT WORKSHOP - 2026  -  DAY 2 - OFFSITE SOLUTIONS DEEP-DIVE"],
        font=TITLE_FONT, size=14, color=TEAL, bold=True)
textbox(s, Inches(0.8), Inches(2.0), Inches(11.7), Inches(1.9),
        ["Understanding Vietnam's DPPA Mechanisms", "and Pricing Considerations"],
        font=TITLE_FONT, size=40, color=NAVY, bold=True, line_spacing=1.05)
textbox(s, Inches(0.8), Inches(4.0), Inches(11.7), Inches(0.9),
        ["The pricing stack - the financial model - what to expect - key decision points",
         "Virtual (grid CfD) DPPA under Decree 57/2025/ND-CP"],
        size=17, color=GREY, line_spacing=1.2)
textbox(s, Inches(0.8), Inches(6.3), Inches(11.7), Inches(0.5),
        ["Allotrope  -  45-minute technical session + Q&A"],
        font=TITLE_FONT, size=13, color=INK)

# ===========================================================================
# SLIDE 2 - Roadmap & objectives
# ===========================================================================
s = content_slide("SESSION ROADMAP", "Follow the money: from your EVN bill to a bankable DPPA")
add_table(s, Inches(0.55), Inches(1.7), Inches(6.7),
          ["#", "Module", "Min"],
          [
              ("1", "The baseline: what you pay EVN today", "4"),
              ("2", "Virtual DPPA money flow: the five-line bill", "12"),
              ("3", "CfD mechanics: what drives your number", "6"),
              ("4", "Developer economics and the three gates", "8"),
              ("5", "Case study: one factory, Cases 5 & 6", "12"),
              ("6", "Key decision points and close", "3"),
          ],
          col_widths=[0.7, 8.0, 1.0], font_size=13, row_h=Inches(0.42))
textbox(s, Inches(7.7), Inches(1.66), Inches(5.1), Inches(0.4),
        ["By the end of this session you can:"], font=TITLE_FONT, size=14,
        color=NAVY, bold=True)
textbox(s, Inches(7.7), Inches(2.1), Inches(5.1), Inches(4.0),
        [
            "1.  Decompose a virtual DPPA monthly settlement into its five components and compare it to your EVN bill",
            "2.  Explain how the CfD strike, market price (FMP), loss factors and grid fees set your net cost",
            "3.  See why buyer, seller and lender each have a gate - and why a deal must pass all three",
            "4.  Leave with realistic pricing expectations grounded in a real factory model",
        ], size=13, line_spacing=1.15, space_after=10)
textbox(s, Inches(0.55), Inches(4.6), Inches(6.7), Inches(0.8),
        ["Assumed from this morning's session: DPPA overview, eligibility, contracting terms.",
         "Scope: virtual (grid CfD) DPPA only."], size=11.5, color=GREY,
        line_spacing=1.15)

# ===========================================================================
# SLIDE 3 - M1: EVN TOU baseline
# ===========================================================================
s = content_slide("MODULE 1 - THE BASELINE (4 MIN)",
                  "What you pay EVN today: time-of-use is the reference price")
add_table(s, Inches(0.55), Inches(1.7), Inches(7.3),
          ["Period", "Hours", "vs base"],
          TOU_PERIODS, col_widths=[2.6, 4.6, 1.2], font_size=12)
textbox(s, Inches(0.55), Inches(3.25), Inches(7.3), Inches(0.35),
        ["Sunday rule: no peak period on Sundays (all hours normal)."],
        size=11.5, color=GREY)
add_table(s, Inches(0.55), Inches(3.75), Inches(7.3),
          ["Voltage", "Off-peak", "Normal", "Peak  (USD/kWh approx. 2025)"],
          TOU_RATES, col_widths=[1.6, 1.3, 1.3, 2.7], font_size=12)
textbox(s, Inches(8.3), Inches(1.66), Inches(4.5), Inches(0.4),
        ["Why this matters for DPPA"], font=TITLE_FONT, size=14, color=NAVY,
        bold=True)
textbox(s, Inches(8.3), Inches(2.1), Inches(4.5), Inches(3.4),
        [
            "Every DPPA offer is judged against this bill - your business-as-usual (BAU)",
            "Peak hours cost ~3.4x off-peak: your load shape decides what \"average\" means for you",
            "Residual energy (load not covered by the DPPA) keeps paying these TOU rates",
        ], size=13, line_spacing=1.15, space_after=10)

# ===========================================================================
# SLIDE 4 - M1: benchmark + eligibility
# ===========================================================================
s = content_slide("MODULE 1 - THE BASELINE (4 MIN)",
                  "The benchmark: EVN BAU with escalation - and who can play")
textbox(s, Inches(0.55), Inches(1.66), Inches(5.9), Inches(0.4),
        ["The savings benchmark"], font=TITLE_FONT, size=14, color=NAVY, bold=True)
textbox(s, Inches(0.55), Inches(2.1), Inches(5.9), Inches(2.6),
        [
            "BAU = your current EVN bill escalated ~4%/yr (historical trend 2015-2024)",
            "A DPPA is evaluated on cumulative cost vs BAU over 10 years and over the project lifetime - not just Year 1",
            "EVN tariff escalation is the main source of long-term DPPA value",
        ], size=13, line_spacing=1.15, space_after=10)
textbox(s, Inches(6.9), Inches(1.66), Inches(5.9), Inches(0.4),
        ["Voltage eligibility (Decree 57/2025, Art. 16)"], font=TITLE_FONT,
        size=14, color=NAVY, bold=True)
add_table(s, Inches(6.9), Inches(2.1), Inches(5.9),
          ["Voltage", "K_pp (2025)", "Loss rate", "Virtual DPPA"],
          KPP_ROWS, col_widths=[1.4, 1.9, 1.0, 1.6], font_size=11.5,
          highlight_rows=(2, 3), highlight_fill=RGBColor(0xFB, 0xEE, 0xDD))
callout(s, Inches(0.55), Inches(5.3), Inches(12.25), Inches(0.9),
        ["If your connection is 6-22kV or below, virtual DPPA is not yet available to you: settlement loss factors (K_pp) are still pending from NLDC/EVN."],
        fill=AMBER, size=13.5)

# ===========================================================================
# SLIDE 5 - M2: money flow diagram
# ===========================================================================
s = content_slide("MODULE 2 - THE FIVE-LINE BILL (12 MIN)",
                  "Virtual DPPA money flow: three relationships, one bill")
# boxes
gen = flow_box(s, Inches(0.7), Inches(2.3), Inches(2.7), Inches(1.25),
               "RE GENERATOR", ["sells 100% of output", "to the spot market"], TEAL)
mkt = flow_box(s, Inches(5.3), Inches(1.75), Inches(2.7), Inches(1.1),
               "VWEM SPOT MARKET", ["pays generator FMP", "per kWh generated"], NAVY)
evn = flow_box(s, Inches(5.3), Inches(3.35), Inches(2.7), Inches(1.1),
               "EVN", ["meters, wheels & bills", "the DPPA customer"], NAVY)
fac = flow_box(s, Inches(9.9), Inches(2.3), Inches(2.7), Inches(1.25),
               "YOUR FACTORY", ["pays EVN the DPPA bill;", "settles CfD with generator"], GREEN)
connector(s, Inches(3.4), Inches(2.6), Inches(5.3), Inches(2.25))
arrow_label(s, Inches(3.35), Inches(1.95), Inches(2.0), "energy -> market (FMP)")
connector(s, Inches(8.0), Inches(3.9), Inches(9.9), Inches(3.3))
arrow_label(s, Inches(7.95), Inches(3.95), Inches(2.1), "DPPA bill (5 lines)")
connector(s, Inches(3.4), Inches(3.2), Inches(9.9), Inches(3.2), color=AMBER,
          dashed=True)
arrow_label(s, Inches(5.6), Inches(2.85), Inches(2.2),
            "CfD settlement (two-way)", color=AMBER)
textbox(s, Inches(0.7), Inches(4.95), Inches(12.1), Inches(1.6),
        [
            "Physical electrons flow through the national grid as usual - the DPPA is a financial overlay",
            "The generator's revenue = spot sales at FMP + CfD settlement with you",
            "Your cost = the five-line EVN bill + the CfD settlement  ->  next slide",
        ], size=13.5, line_spacing=1.2, space_after=8)

# ===========================================================================
# SLIDE 6 - M2: five-line bill table
# ===========================================================================
s = content_slide("MODULE 2 - THE FIVE-LINE BILL (12 MIN)",
                  "Your monthly settlement has five lines")
add_table(s, Inches(0.55), Inches(1.7), Inches(12.25),
          ["#", "Line", "Formula", "Level (2025 reference)"],
          [
              ("1", "Market energy", "Q_Khc x FMP x k x K_pp",
               "FMP avg ~" + FMP_2025 + " VND/kWh; k~1.026; K_pp by voltage"),
              ("2", "DPPA service fee", "Q_Khc x C_dppa_dv", "360 VND/kWh"),
              ("3", "Balancing fee", "Q_Khc x P_cl", "163.3 VND/kWh"),
              ("4", "Residual purchase", "(Load - Q_Khc) x EVN retail",
               "your TOU rates (Module 1)"),
              ("5", "CfD settlement", "(Strike - FMP) x min(Q_c, Q_Khc)",
               "two-way: can be + or -"),
          ],
          col_widths=[0.5, 2.2, 4.2, 5.3], font_size=12.5, row_h=Inches(0.52),
          highlight_rows=(4,), highlight_fill=RGBColor(0xFB, 0xEE, 0xDD))
textbox(s, Inches(0.55), Inches(5.05), Inches(12.25), Inches(1.6),
        [
            "Lines 1-4 are billed by EVN. Line 5 settles bilaterally with the generator under your CfD contract.",
            "Q_Khc = the renewable volume you actually received and consumed (defined next slide).",
            "CFMP = FMP x k: k converts the market price to the customer delivery point (price conversion, not a volume loss).",
        ], size=12.5, line_spacing=1.2, space_after=6)

# ===========================================================================
# SLIDE 7 - M2: volume definitions
# ===========================================================================
s = content_slide("MODULE 2 - THE FIVE-LINE BILL (12 MIN)",
                  "Which volume? You settle only on what you actually consume")
formula_box(s, Inches(0.55), Inches(1.75), Inches(7.4), Inches(2.0),
            [
                "Q_adj = Q_meter / K_pp        delivered volume at your meter",
                "Q_Khc = min(Load, Q_adj)      settled volume (lines 1-3)",
                "Q_CfD = min(Q_c, Q_Khc)       CfD settles on received volume",
            ], size=14)
add_table(s, Inches(8.3), Inches(1.75), Inches(4.5),
          ["Voltage", "K_pp", "Loss"],
          [(r[0], r[1], r[2]) for r in KPP_ROWS[:2]],
          col_widths=[1.5, 1.8, 1.0], font_size=12)
textbox(s, Inches(0.55), Inches(4.1), Inches(12.25), Inches(2.2),
        [
            "K_pp removes grid losses: ~2.8-4.6% of plant output never reaches your meter - you do not pay for it, but you also do not receive it",
            "If the plant generates more than you consume in an hour, the excess earns the generator spot revenue only - it is NOT billed to you and does NOT earn CfD",
            "Practical consequence: contracted volume (Q_c) above your real consumption buys nothing",
        ], size=13, line_spacing=1.2, space_after=10)

# ===========================================================================
# SLIDE 8 - M2: worked example inputs (NSMO Vi du 1)
# ===========================================================================
s = content_slide("MODULE 2 - WORKED EXAMPLE (NSMO, SEP-2025)",
                  "Official NSMO simulation: industrial park buys 6,000 MWh/month")
add_table(s, Inches(0.55), Inches(1.7), Inches(6.6),
          ["Parameter", "Symbol", "Value"],
          [
              ("Consumption = received RE volume", "Q_KH = Q_Khc", "6,000,000 kWh/mo"),
              ("Contracted CfD volume", "Q_c", "6,000,000 kWh/mo"),
              ("CfD strike price", "P_c", VD1["pc"] + " VND/kWh"),
              ("Market price (monthly avg)", "FMP", VD1["fmp"] + " VND/kWh"),
              ("Price conversion factor", "k", VD1["k"]),
              ("Loss factor (110 kV)", "K_pp", VD1["kpp"]),
              ("DPPA service fee", "C_dppa_dv", VD1["cdppa"] + " VND/kWh"),
              ("Balancing fee", "P_cl", VD1["pcl"] + " VND/kWh"),
              ("Avg retail price (residual)", "P1", VD1["p1"] + " VND/kWh"),
          ],
          col_widths=[3.6, 1.5, 2.2], font_size=12, row_h=Inches(0.38))
textbox(s, Inches(7.6), Inches(1.66), Inches(5.2), Inches(0.4),
        ["Scenario"], font=TITLE_FONT, size=14, color=NAVY, bold=True)
textbox(s, Inches(7.6), Inches(2.1), Inches(5.2), Inches(3.6),
        [
            "A 50 MW solar plant contracts its full output to one industrial customer",
            "Consumption exactly equals delivered RE volume - no residual EVN purchase this month",
            "Source: NSMO worked examples for Decree 57/2025 (Sep-2025 public training)",
            "Strike just 100 VND above market - watch what the fees do",
        ], size=13, line_spacing=1.15, space_after=10)

# ===========================================================================
# SLIDE 9 - M2: worked example arithmetic
# ===========================================================================
s = content_slide("MODULE 2 - WORKED EXAMPLE (NSMO, SEP-2025)",
                  "The five lines, computed")
add_table(s, Inches(0.55), Inches(1.7), Inches(12.25),
          ["Line", "Calculation", "VND/month"],
          [
              ("1  Market energy (C_DN)",
               "6,000,000 x 1,200 x 1.026 x 1.008", "7,446,297,600"),
              ("2  Service fee (C_dppa)", "6,000,000 x 360", "2,160,000,000"),
              ("3  Balancing fee (C_cl)", "6,000,000 x 163.3", "979,800,000"),
              ("4  Residual (C_bl)", "0 kWh x 2,204", "0"),
              ("   EVN bill (C_EVN)", "lines 1 + 2 + 3 + 4",
               ("10,586,097,600", {"bold": True})),
              ("5  CfD settlement (C_CfD)", "(1,300 - 1,200) x 6,000,000",
               "600,000,000"),
              ("   TOTAL (C_KH)", "C_EVN + C_CfD",
               ("11,186,097,600", {"bold": True, "color": NAVY})),
          ],
          col_widths=[3.0, 4.6, 2.6], font_size=12.5, row_h=Inches(0.42),
          highlight_rows=(4, 6))
textbox(s, Inches(0.55), Inches(5.15), Inches(12.25), Inches(1.4),
        [
            "Effective cost = 11,186,097,600 / 6,000,000 = ~1,864 VND/kWh for fully renewable supply (vs ~2,204 VND avg retail in the example)",
            "Generator side: 6,048,000 kWh x 1,200 (spot) + 600,000,000 (CfD) = 7,857,600,000 VND",
        ], size=12.5, line_spacing=1.2, space_after=6)

# ===========================================================================
# SLIDE 10 - M2: grid fee punchline
# ===========================================================================
s = content_slide("MODULE 2 - THE FIVE-LINE BILL (12 MIN)",
                  "The structural cost: 523 VND/kWh of fees before any CfD")
add_table(s, Inches(0.55), Inches(1.75), Inches(7.6),
          ["Component (per matched kWh)", "VND/kWh", "US cents"],
          [
              ("Market energy at 2025 avg FMP x k x K_pp (22-110kV)", "~1,504", "~6.0c"),
              ("DPPA service fee (C_dppa_dv)", "360", "1.4c"),
              ("Balancing fee (P_cl)", "163.3", "0.7c"),
              ("Delivered cost before CfD", ("~2,027", {"bold": True}),
               ("~8.1c", {"bold": True})),
          ],
          col_widths=[4.6, 1.4, 1.2], font_size=12.5, row_h=Inches(0.46),
          highlight_rows=(3,))
textbox(s, Inches(8.6), Inches(1.75), Inches(4.2), Inches(3.4),
        [
            "The 523 VND/kWh of fixed fees applies to every matched kWh, regardless of your strike",
            "They are set by regulation, not negotiation",
            "Compare ~2,027 VND delivered (pre-CfD) with your own TOU mix - for many factories that is already near or above average grid cost",
        ], size=13, line_spacing=1.15, space_after=10)
callout(s, Inches(0.55), Inches(5.5), Inches(12.25), Inches(1.0),
        ["Virtual DPPA is rarely a Day-1 discount. Its value = price certainty vs EVN escalation + credible renewable claims. Treat fee levels as the first number you check."],
        size=14)

# ===========================================================================
# SLIDE 11 - M3: two-way CfD
# ===========================================================================
s = content_slide("MODULE 3 - CFD MECHANICS (6 MIN)",
                  "The CfD is two-way: certainty, not automatic savings")
formula_box(s, Inches(0.55), Inches(1.75), Inches(6.6), Inches(0.85),
            ["CfD cash flow = (Strike - FMP) x Q_CfD"], size=16)
add_table(s, Inches(0.55), Inches(2.95), Inches(6.6),
          ["Market condition", "Who pays whom"],
          [
              ("FMP below strike", "You top up the generator"),
              ("FMP above strike", "Generator pays you back"),
          ],
          col_widths=[2.6, 4.0], font_size=13, row_h=Inches(0.46))
textbox(s, Inches(7.7), Inches(1.75), Inches(5.1), Inches(4.2),
        [
            "Net effect: your renewable energy cost is locked near the strike (+ fees), whatever the market does",
            "2025 reference: FMP averaged ~1,427 VND/kWh with strong hourly volatility - the CfD removes that volatility from your budget",
            "A strike below expected average FMP would mean expected receipts - but no developer can finance that (Module 4)",
        ], size=13, line_spacing=1.2, space_after=12)

# ===========================================================================
# SLIDE 12 - M3: escalation & Year 1
# ===========================================================================
s = content_slide("MODULE 3 - CFD MECHANICS (6 MIN)",
                  "Strike escalation and the Year 1 reality")
textbox(s, Inches(0.55), Inches(1.7), Inches(6.0), Inches(0.4),
        ["Strike escalation (e.g. 4%/yr)"], font=TITLE_FONT, size=14,
        color=NAVY, bold=True)
textbox(s, Inches(0.55), Inches(2.12), Inches(6.0), Inches(2.6),
        [
            "Escalating strikes improve the seller's and lender's lifetime economics...",
            "...but do nothing for your Year 1 cost - and compound against you later",
            "Negotiate the escalation index as hard as the starting strike (fixed VND vs CPI vs USD-linked)",
        ], size=13, line_spacing=1.18, space_after=10)
textbox(s, Inches(7.0), Inches(1.7), Inches(5.8), Inches(0.4),
        ["What Year 1 usually looks like"], font=TITLE_FONT, size=14,
        color=NAVY, bold=True)
textbox(s, Inches(7.0), Inches(2.12), Inches(5.8), Inches(2.6),
        [
            "In most realistic structures, Year 1 DPPA cost is at or above BAU",
            "Savings build as EVN tariffs escalate (~4%/yr) against your locked strike",
            "Evaluate on 10-year and lifetime cumulative cost - never on Year 1 alone",
        ], size=13, line_spacing=1.18, space_after=10)
callout(s, Inches(0.55), Inches(5.2), Inches(12.25), Inches(0.9),
        ["Budget question for your CFO: can we accept a modest Year 1 premium to cap a 25-year exposure?"],
        fill=TEAL, size=14)

# ===========================================================================
# SLIDE 13 - M3: volume vs load shape
# ===========================================================================
s = content_slide("MODULE 3 - CFD MECHANICS (6 MIN)",
                  "Contract volume vs your load shape")
textbox(s, Inches(0.55), Inches(1.7), Inches(12.25), Inches(0.7),
        ["Solar generates midday; your factory runs its own shape. The overlap - hour by hour - is all that counts."],
        size=14, color=INK)
add_table(s, Inches(0.55), Inches(2.5), Inches(12.25),
          ["If you...", "What happens", "Cost consequence"],
          [
              ("Over-contract (Q_c > consumption)",
               "CfD caps at consumed volume - excess earns nothing",
               "You paid contract terms for volume you can never settle"),
              ("Match contract to hourly overlap",
               "Every contracted kWh settles", "Efficient - the target"),
              ("Under-contract (Q_c < overlap)",
               "More residual purchased from EVN at TOU",
               "Less price certainty, smaller renewable claim"),
          ],
          col_widths=[3.4, 4.6, 4.2], font_size=12.5, row_h=Inches(0.62),
          highlight_rows=(1,))
textbox(s, Inches(0.55), Inches(5.3), Inches(12.25), Inches(1.2),
        [
            "Ask every developer for an hourly (8760) matching analysis against YOUR meter data - not an annual percentage",
            "Batteries can shift the overlap window but add capex the tariff must carry (see the case study)",
        ], size=13, line_spacing=1.2, space_after=8)

# ===========================================================================
# SLIDE 14 - M4: developer revenue
# ===========================================================================
s = content_slide("MODULE 4 - DEVELOPER ECONOMICS (8 MIN)",
                  "The other side of the table: where the developer's money comes from")
formula_box(s, Inches(0.55), Inches(1.75), Inches(7.2), Inches(1.3),
            [
                "Revenue = Q_meter x FMP          spot sales (volatile)",
                "        + (Strike - FMP) x Q_CfD  CfD with you (stabiliser)",
            ], size=14)
textbox(s, Inches(8.3), Inches(1.75), Inches(4.5), Inches(3.2),
        [
            "The CfD converts volatile spot revenue into a near-fixed price on matched volume",
            "Your strike IS the developer's bankable revenue line",
            "Everything the lender checks flows from it",
        ], size=13, line_spacing=1.2, space_after=12)
textbox(s, Inches(0.55), Inches(3.5), Inches(7.2), Inches(2.8),
        [
            "Typical Vietnam RE project finance:",
            ("  70% debt / 30% equity, ~10-year debt term"),
            ("  Debt: ~8.5%/yr in VND  -  ~5%/yr in USD (international)"),
            ("  Equity target: 12-15%+ IRR"),
            ("  CIT incentives: 4 years exempt + 9 years at half rate"),
        ], size=13, line_spacing=1.2, space_after=6)

# ===========================================================================
# SLIDE 15 - M4: the three gates
# ===========================================================================
s = content_slide("MODULE 4 - DEVELOPER ECONOMICS (8 MIN)",
                  "A deal must pass three gates at once")
add_table(s, Inches(0.55), Inches(1.75), Inches(12.25),
          ["Gate", "Who", "Test", "Strike pressure"],
          [
              ("BUYER", "You",
               "Cumulative cost <= BAU over 10 yr AND lifetime", "pushes strike DOWN"),
              ("SELLER", "Developer", "Equity IRR >= 12-15%", "pushes strike UP"),
              ("LENDER", "Bank", "Minimum DSCR >= ~1.20x every year",
               "pushes strike UP (hardest gate)"),
          ],
          col_widths=[1.5, 1.5, 5.6, 3.2], font_size=13, row_h=Inches(0.55))
textbox(s, Inches(0.55), Inches(4.0), Inches(12.25), Inches(1.6),
        [
            "DSCR = cash available for debt service / debt service, computed for every year of the loan",
            "One bad year - e.g. a battery replacement - can sink the whole financing even if averages look fine",
            "The negotiation window is the strike range where all three gates pass - it can be empty (case study)",
        ], size=13, line_spacing=1.2, space_after=8)
callout(s, Inches(0.55), Inches(5.7), Inches(12.25), Inches(0.85),
        ["A strike below the developer's bankability floor does not mean a cheap deal - it means no project."],
        size=14)

# ===========================================================================
# SLIDE 16 - M4: lifecycle items
# ===========================================================================
s = content_slide("MODULE 4 - DEVELOPER ECONOMICS (8 MIN)",
                  "Lifecycle items the tariff must carry - ask about all four")
add_table(s, Inches(0.55), Inches(1.75), Inches(12.25),
          ["Item", "Typical assumption", "Why you should care"],
          [
              ("PV degradation", "~0.5%/yr output decline",
               "Matched volume shrinks; lost kWh repurchased from EVN at TOU"),
              ("Battery replacement", "Around year 10-11, large one-off cost",
               "Can crater DSCR in that year; ask who funds it and how"),
              ("CIT holiday", "4 yr exempt + 9 yr half rate (RE projects)",
               "Improves developer economics - value you can negotiate for"),
              ("FX exposure", "VND strike vs USD capex/debt",
               "USD-linked terms shift FX risk to you; price it consciously"),
          ],
          col_widths=[2.4, 4.0, 5.8], font_size=12.5, row_h=Inches(0.6))
textbox(s, Inches(0.55), Inches(4.8), Inches(12.25), Inches(1.2),
        ["These rarely appear in headline term sheets, but each one moves the bankable strike floor - and therefore your price."],
        size=13, color=GREY)

# ===========================================================================
# SLIDE 17 - M5: case study framing
# ===========================================================================
s = content_slide("MODULE 5 - CASE STUDY (12 MIN)",
                  "Same factory as yesterday - now going off-site")
textbox(s, Inches(0.55), Inches(1.7), Inches(12.25), Inches(0.7),
        ["Yesterday (On-Site session): Cases 1-4 explored rooftop solar + BESS behind the meter. Today: the same southern-Vietnam manufacturing facility (22-110kV) signs a virtual DPPA with a dedicated solar + storage project."],
        size=13.5, line_spacing=1.2)
add_table(s, Inches(0.55), Inches(2.75), Inches(12.25),
          ["Deal frame (both cases)", "Value"],
          [
              ("DPPA type", "Virtual (grid CfD) under Decree 57/2025"),
              ("Strike price", "2,000 VND/kWh, escalating 4%/yr"),
              ("Tenor / analysis period", "25 years"),
              ("Financing", "70% debt, 8.5% VND, 10-yr term"),
              ("Difference between cases",
               ("Case 5: solar + LARGE battery   -   Case 6: solar + MINIMUM battery",
                {"bold": True})),
          ],
          col_widths=[3.6, 8.6], font_size=13, row_h=Inches(0.46))
textbox(s, Inches(0.55), Inches(5.75), Inches(12.25), Inches(0.8),
        ["Method: full 8760-hour settlement simulation (ND57 formulas, hourly FMP/CFMP series) + 25-year project cash-flow model with the three gates."],
        size=12, color=GREY)

# ===========================================================================
# SLIDE 18 - M5: Case 5
# ===========================================================================
s = content_slide("MODULE 5 - CASE STUDY (12 MIN)",
                  "Case 5 - solar + large battery: the battery eats the deal")
add_table(s, Inches(0.55), Inches(1.75), Inches(6.4),
          ["Metric", "Result"],
          [
              ("Seller equity IRR", CASE5["irr"]),
              ("Project IRR", CASE5["pirr"]),
              ("Developer NPV", CASE5["npv"]),
              ("Minimum DSCR", (CASE5["dscr"] + "  (replacement year)", {"color": RED})),
              ("Payback", CASE5["payback"]),
              ("Buyer vs BAU - Year 1", (CASE5["b1"], {"color": RED})),
              ("Buyer vs BAU - 10-yr cumulative", (CASE5["b10"], {"color": RED})),
              ("Buyer vs BAU - lifetime", (CASE5["blife"], {"color": RED})),
          ],
          col_widths=[3.6, 2.8], font_size=12.5, row_h=Inches(0.42))
textbox(s, Inches(7.5), Inches(1.75), Inches(5.3), Inches(3.6),
        [
            "At strike 2,000 the buyer pays ~9% MORE than BAU on every horizon",
            "The oversized battery adds a ~$1.2M replacement around year 11 that the tariff must recover",
            "Even so, lender comfort is thin (min DSCR 1.14x < 1.20x)",
        ], size=13, line_spacing=1.2, space_after=12)
callout(s, Inches(0.55), Inches(5.6), Inches(12.25), Inches(0.9),
        ["Lesson 1: an oversized battery cannot carry its own replacement cost through the tariff - everyone loses."],
        fill=RED, size=14)

# ===========================================================================
# SLIDE 19 - M5: Case 6
# ===========================================================================
s = content_slide("MODULE 5 - CASE STUDY (12 MIN)",
                  "Case 6 - solar + minimum battery: bankable, but the strike is the fight")
add_table(s, Inches(0.55), Inches(1.75), Inches(6.4),
          ["Metric", "Result"],
          [
              ("Seller equity IRR", (CASE6["irr"], {"color": GREEN})),
              ("Project IRR", CASE6["pirr"]),
              ("Developer NPV", CASE6["npv"]),
              ("Minimum DSCR", (CASE6["dscr"] + "  (comfortably > 1.20x)", {"color": GREEN})),
              ("Payback", CASE6["payback"]),
              ("Buyer vs BAU - all horizons", (CASE6["blife"], {"color": RED})),
          ],
          col_widths=[3.6, 2.8], font_size=12.5, row_h=Inches(0.46))
textbox(s, Inches(7.5), Inches(1.75), Inches(5.3), Inches(3.6),
        [
            "Lean sizing: seller and lender are now comfortable - returns strong, DSCR 1.50x",
            "But at strike 2,000 the buyer is ~14% worse than BAU: the seller's headroom IS the negotiation space",
            "So: sweep the strike down and see where the gates cross",
        ], size=13, line_spacing=1.2, space_after=12)
callout(s, Inches(0.55), Inches(5.6), Inches(12.25), Inches(0.9),
        ["Lesson 2: right-sized assets make the project financeable - then the entire negotiation is the strike price."],
        fill=TEAL, size=14)

# ===========================================================================
# SLIDE 20 - M5: negotiation sweep
# ===========================================================================
s = content_slide("MODULE 5 - CASE STUDY (12 MIN)",
                  "Sweeping the strike on Case 6: the window is empty")
textbox(s, Inches(0.55), Inches(1.7), Inches(12.25), Inches(0.45),
        ["56 scenarios swept: strikes 1,200-2,200 VND/kWh x contract volume 70-100%. Gate results:"],
        size=13)
add_table(s, Inches(0.55), Inches(2.25), Inches(12.25),
          ["Strike level", "Buyer gate", "Seller gate", "Lender gate (>=1.20x)", "Balanced?"],
          [
              ("~2,000 (offer)", ("FAIL  -14%", {"color": RED}), "PASS", "PASS  1.50x", "No"),
              ("~1,400", ("FAIL  -1.4%", {"color": RED}), "PASS  ~19%", "PASS  1.19-1.5x", "No"),
              ("~1,300 x 70% vol", ("PASS  +0.5% lifetime", {"color": GREEN}),
               "PASS  17.9%", ("FAIL  1.14x", {"color": RED}), "No"),
              ("~1,200", ("PASS  +2.9%", {"color": GREEN}), "PASS",
               ("FAIL  <1.20x", {"color": RED}), "No"),
          ],
          col_widths=[2.4, 2.5, 2.0, 2.9, 1.4], font_size=12.5, row_h=Inches(0.5))
textbox(s, Inches(0.55), Inches(4.65), Inches(12.25), Inches(0.9),
        ["The buyer turns positive just as the lender drops out: zero of 56 scenarios pass all three gates at current market prices and fee levels."],
        size=13.5, bold=True, color=NAVY)
callout(s, Inches(0.55), Inches(5.6), Inches(12.25), Inches(0.95),
        ["Pricing alone cannot balance this deal. The levers that can: lower leverage, debt sculpting, USD-denominated debt (~5%), fee relief, longer tenor - exactly this afternoon's panel."],
        fill=AMBER, size=13.5)

# ===========================================================================
# SLIDE 21 - M5: takeaways
# ===========================================================================
s = content_slide("MODULE 5 - CASE STUDY (12 MIN)",
                  "What this means for your pricing expectations")
textbox(s, Inches(0.55), Inches(1.8), Inches(12.25), Inches(4.4),
        [
            ("1.  Expect low single-digit lifetime savings at best - the product you are buying is 25-year price certainty plus credible renewable claims, not a discount", {"size": 15}),
            ("2.  Expect Year 1 to cost more than BAU - plan the internal approval around cumulative horizons", {"size": 15}),
            ("3.  Strike offers near 2,000 VND/kWh leave the buyer deeply negative - realistic buyer-positive territory starts near ~1,300, which today fails the lender", {"size": 15}),
            ("4.  Deal feasibility is decided by financing structure as much as by price - ask developers HOW their project is financed, not just what they charge", {"size": 15}),
        ], line_spacing=1.25, space_after=16)

# ===========================================================================
# SLIDE 22 - M6: decision checklist
# ===========================================================================
s = content_slide("MODULE 6 - DECISION POINTS (3 MIN)",
                  "Five decisions to settle before you sign")
add_table(s, Inches(0.55), Inches(1.75), Inches(12.25),
          ["#", "Decision", "What to check"],
          [
              ("1", "Strike level & escalation/indexation",
               "Fixed VND vs escalating vs USD-linked; model all three vs BAU"),
              ("2", "Contract volume vs load shape",
               "8760 hourly matching vs YOUR meter data; CfD caps at consumed volume"),
              ("3", "Tenor & termination",
               "Developer wants 15-25 yr; match to your lease / order book; exit terms"),
              ("4", "Voltage eligibility & K_pp",
               ">=22kV only today; which K_pp applies decides your loss adjustment"),
              ("5", "Risk allocation",
               "FMP volatility, curtailment, negative-CfD months, change-in-law -> this afternoon"),
          ],
          col_widths=[0.6, 4.0, 7.6], font_size=12.5, row_h=Inches(0.56))
textbox(s, Inches(0.55), Inches(5.4), Inches(12.25), Inches(0.6),
        ["Items 1-4 are quantifiable before negotiation - run them against your own load data first."],
        size=12.5, color=GREY)

# ===========================================================================
# SLIDE 23 - Close
# ===========================================================================
s = content_slide("CLOSE", "Three things to remember")
textbox(s, Inches(0.55), Inches(1.9), Inches(12.25), Inches(3.4),
        [
            ("1.  Know your five-line bill - 523 VND/kWh of fees and the loss factor are fixed; the strike is the only big number you negotiate", {"size": 16}),
            ("2.  Judge the deal on 10-year and lifetime cost vs BAU - never on Year 1", {"size": 16}),
            ("3.  Every deal must pass buyer, seller AND lender gates - if the window is empty, the fix is structural, not just price", {"size": 16}),
        ], line_spacing=1.3, space_after=18)
callout(s, Inches(0.55), Inches(5.4), Inches(12.25), Inches(1.0),
        ["Next: panel discussion on risk allocation and the reality of bank financing for synthetic DPPA - bring the lender-gate question with you."],
        fill=NAVY, size=14)

# ===========================================================================
# BACKUP SLIDES
# ===========================================================================
s = content_slide("BACKUP", "Backup slides - Q&A reference", accent=GREY)
textbox(s, Inches(0.55), Inches(2.2), Inches(12.25), Inches(2.4),
        [
            "B1   FMP market context and forward uncertainty",
            "B2   What happens if my load drops?",
            "B3   RECs and Scope 2 treatment of DPPA volumes",
            "B4   6-22kV connections: current status",
            "B5   How the buyer gate is measured",
        ], size=15, line_spacing=1.3, space_after=10)

# --- B1
s = content_slide("BACKUP B1", "FMP market context (2025)", accent=GREY)
add_table(s, Inches(0.55), Inches(1.75), Inches(6.6),
          ["Price", "VND/kWh", "USD/kWh (25,000)"],
          [
              ("FMP, 2025 mean", FMP_2025, "0.0571"),
              ("CFMP, 2025 mean", CFMP_2025, "0.0586"),
              ("Measured k ratio (hourly CFMP/FMP)", "1.006 - 1.060", "-"),
          ],
          col_widths=[3.2, 1.7, 1.7], font_size=12.5, row_h=Inches(0.46))
textbox(s, Inches(0.55), Inches(3.6), Inches(12.25), Inches(2.2),
        [
            "FMP is strongly seasonal and hourly-volatile; monthly averages mask large swings",
            "Forward FMP is the single biggest uncertainty in any DPPA business case - run high/low scenarios",
            "Figures are 2025 VWEM data - refresh before relying on them for a 2026+ decision",
        ], size=13, line_spacing=1.2, space_after=10)

# --- B2
s = content_slide("BACKUP B2", "What happens if my load drops?", accent=GREY)
textbox(s, Inches(0.55), Inches(1.8), Inches(12.25), Inches(3.6),
        [
            "Settlement shrinks with you: Q_Khc = min(Load, Q_adj) - EVN lines 1-3 only ever bill consumed volume",
            "The CfD also caps at consumed volume - but your CONTRACT may still oblige payment on contracted volume (take-or-pay style clauses)",
            "This is a contract negotiation point, not a regulatory mechanic: check shortfall, suspension and force-majeure clauses",
            "Match tenor and volume commitments to your realistic production outlook, not the optimistic one",
        ], size=13.5, line_spacing=1.25, space_after=12)

# --- B3
s = content_slide("BACKUP B3", "RECs and Scope 2 treatment", accent=GREY)
textbox(s, Inches(0.55), Inches(1.8), Inches(12.25), Inches(3.6),
        [
            "DPPA volumes carry the renewable attributes you need for market-based Scope 2 accounting (see Day 1, Topic 2)",
            "Confirm in the contract WHO owns the attributes/RECs for matched volume - and for the plant's surplus sold to market",
            "Hourly settlement data (Q_Khc) gives you credible granular matching evidence for brand reporting",
        ], size=13.5, line_spacing=1.25, space_after=12)

# --- B4
s = content_slide("BACKUP B4", "6-22kV connections: current status", accent=GREY)
textbox(s, Inches(0.55), Inches(1.8), Inches(12.25), Inches(3.6),
        [
            "Decree 57/2025 (Art. 16) enables virtual DPPA for >=110kV and 22-110kV connections",
            "K_pp loss factors for 6-22kV and <6kV have not been published by NLDC/EVN - settlement cannot be computed",
            "Practical options today: connect/upgrade to >=22kV, pursue on-site solutions (Day 1), or wait for the loss-factor announcement",
            "Many supplier factories sit at 6-22kV - raise this with CEBA and your brand contacts as a policy ask",
        ], size=13.5, line_spacing=1.25, space_after=12)

# --- B5
s = content_slide("BACKUP B5", "How the buyer gate is measured", accent=GREY)
textbox(s, Inches(0.55), Inches(1.8), Inches(12.25), Inches(3.8),
        [
            "Benchmark: EVN-only BAU (your current tariff escalated ~4%/yr), same load, same horizon",
            "Gate: cumulative DPPA cost <= cumulative BAU cost over BOTH the first 10 years AND the project lifetime",
            "Year 1 outcome is reported as context but is not the gate",
            "A separate comparison vs an on-site alternative is useful disclosure, but the qualification test is vs BAU",
        ], size=13.5, line_spacing=1.25, space_after=12)

# ---------------------------------------------------------------------------
out = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                   "CEBA_DPPA_Mechanisms_Pricing.pptx")
prs.save(out)
print("Saved", out, "with", len(prs.slides._sldIdLst), "slides")
