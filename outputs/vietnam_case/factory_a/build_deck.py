"""Build the Factory A Solar+BESS case-study deck (7 slides, corporate, white bg)."""
import json
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION, XL_TICK_MARK

M = json.load(open("clean_metrics.json"))

# ---- design system -------------------------------------------------------
INK   = RGBColor(0x0F, 0x29, 0x3D)   # deep navy (headlines)
TEAL  = RGBColor(0x12, 0x7A, 0x7A)   # primary accent (clean energy)
GREEN = RGBColor(0x2E, 0x9E, 0x5B)   # solar
AMBER = RGBColor(0xE0, 0x9A, 0x2B)   # BESS / storage
RED   = RGBColor(0xC0, 0x4A, 0x3A)   # grid / cost
SLATE = RGBColor(0x55, 0x66, 0x77)   # body text
MUTE  = RGBColor(0x8A, 0x98, 0xA5)   # captions
LIGHT = RGBColor(0xF1, 0xF5, 0xF7)   # light fill
CARD  = RGBColor(0xF7, 0xF9, 0xFA)
LINE  = RGBColor(0xDD, 0xE4, 0xE9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
HEAD = "Segoe UI Semibold"
BODY = "Segoe UI"
LIGHTF = "Segoe UI Light"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.fill.solid(); bg.fill.fore_color.rgb = WHITE; bg.line.fill.background()
    bg.shadow.inherit = False
    s.shapes._spTree.remove(bg._element); s.shapes._spTree.insert(2, bg._element)
    return s


def _no_shadow(sp):
    """Remove the inherited preset shadow so shapes render flat.

    The shadow comes from the theme via the shape <p:style>/<a:effectRef>,
    not from an effectLst. Drop the whole style block (fill/line are set
    explicitly) and add an empty effectLst to be safe.
    """
    el = sp._element
    for style in el.findall(qn("p:style")):
        el.remove(style)
    spPr = el.spPr
    for tag in ("a:effectLst", "a:effectDag"):
        for e in spPr.findall(qn(tag)):
            spPr.remove(e)
    spPr.append(spPr.makeelement(qn("a:effectLst"), {}))


def rect(s, x, y, w, h, fill=None, line=None, lw=1.0, shape=MSO_SHAPE.RECTANGLE):
    sp = s.shapes.add_shape(shape, x, y, w, h)
    sp.shadow.inherit = False
    _no_shadow(sp)
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(lw)
    return sp


def txt(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        space_after=2, line_spacing=1.0, wrap=True):
    """runs: list of paragraphs; each paragraph is list of (text, size, color, font, bold)."""
    tb = s.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    for m in ("left", "right", "top", "bottom"):
        setattr(tf, f"margin_{m}", 0)
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space_after); p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for (t, sz, col, fn, bold) in para:
            r = p.add_run(); r.text = t
            r.font.size = Pt(sz); r.font.color.rgb = col
            r.font.name = fn; r.font.bold = bold
    return tb


def header(s, kicker, title, accent=TEAL):
    rect(s, Inches(0.55), Inches(0.52), Inches(0.09), Inches(0.62), fill=accent)
    txt(s, Inches(0.78), Inches(0.48), Inches(11.8), Inches(0.3),
        [[(kicker.upper(), 11, accent, HEAD, True)]])
    txt(s, Inches(0.78), Inches(0.74), Inches(11.9), Inches(0.62),
        [[(title, 25, INK, HEAD, True)]])
    rect(s, Inches(0.78), Inches(1.42), Inches(11.78), Pt(1.4), fill=LINE)


def footer(s, n):
    txt(s, Inches(0.78), Inches(7.05), Inches(8), Inches(0.3),
        [[("Factory A · Solar + BESS Techno-Economic Case Study", 8.5, MUTE, BODY, False)]])
    txt(s, Inches(11.4), Inches(7.05), Inches(1.16), Inches(0.3),
        [[(f"{n:02d}", 8.5, MUTE, HEAD, True)]], align=PP_ALIGN.RIGHT)


def kpi(s, x, y, w, h, value, vcol, label, sub=None, fill=CARD, edge=LINE):
    rect(s, x, y, w, h, fill=fill, line=edge, lw=1.0)
    rect(s, x, y, Inches(0.07), h, fill=vcol)
    paras = [[(value, 25, vcol, HEAD, True)],
             [(label, 11, INK, HEAD, True)]]
    if sub:
        paras.append([(sub, 9.5, SLATE, BODY, False)])
    txt(s, x + Inches(0.26), y + Inches(0.18), w - Inches(0.42), h - Inches(0.3),
        paras, anchor=MSO_ANCHOR.MIDDLE, space_after=3, line_spacing=1.0)


def style_chart(chart, colors, num_fmt="0", lab_size=10, legend=False, max_scale=None):
    chart.has_title = False
    plot = chart.plots[0]
    plot.has_data_labels = True
    dl = plot.data_labels
    dl.number_format = num_fmt; dl.number_format_is_linked = False
    dl.font.size = Pt(lab_size); dl.font.bold = True; dl.font.name = HEAD
    dl.position = XL_LABEL_POSITION.OUTSIDE_END
    plot.gap_width = 70
    cat = chart.category_axis
    cat.tick_labels.font.size = Pt(10.5); cat.tick_labels.font.name = BODY
    cat.tick_labels.font.color.rgb = INK
    cat.format.line.color.rgb = LINE
    cat.major_tick_mark = XL_TICK_MARK.NONE
    cat.minor_tick_mark = XL_TICK_MARK.NONE
    val = chart.value_axis
    val.visible = False
    val.has_major_gridlines = False
    val.minimum_scale = 0
    if max_scale:
        val.maximum_scale = max_scale
    chart.has_legend = legend
    if legend:
        chart.legend.position = XL_LEGEND_POSITION.TOP
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(10); chart.legend.font.name = BODY
    ser = chart.series[0]
    for i, c in enumerate(colors):
        pt = ser.points[i]
        pt.format.fill.solid(); pt.format.fill.fore_color.rgb = c
        pt.format.line.fill.background()
    return chart


SC = ["Solar Only", "Solar + BESS\nCurrent TOU", "Solar + BESS\nTOU 963", "Solar + BESS\n2-Component"]
KEYS = ["case_4", "case_1", "case_2", "case_3"]

# =========================================================================
# SLIDE 1 — TITLE
# =========================================================================
s = slide()
rect(s, 0, 0, SW, SH, fill=WHITE)
# left accent column
rect(s, 0, 0, Inches(0.32), SH, fill=INK)
rect(s, Inches(0.32), 0, Inches(0.06), SH, fill=TEAL)
txt(s, Inches(0.95), Inches(1.15), Inches(11), Inches(0.4),
    [[("CEBA CLEAN ENERGY PROCUREMENT WORKSHOP · 2026", 12.5, TEAL, HEAD, True)]])
txt(s, Inches(0.9), Inches(1.95), Inches(11.6), Inches(2.2),
    [[("Solar + BESS:", 52, INK, LIGHTF, False)],
     [("The Economic Case", 52, INK, HEAD, True)]],
    line_spacing=1.02, space_after=2)
rect(s, Inches(0.95), Inches(4.25), Inches(3.2), Pt(2.4), fill=AMBER)
txt(s, Inches(0.95), Inches(4.5), Inches(11), Inches(0.9),
    [[("A techno-economic analysis of four clean-energy procurement strategies for a ",
       15, SLATE, BODY, False),
      ("single Vietnamese manufacturer — Factory A.", 15, INK, HEAD, True)]],
    line_spacing=1.2)
# bottom method chips
chips = ["REopt MILP system optimization", "25-year ESCO pro forma", "70% project debt",
         "9.3 GWh/yr load · 22–110 kV"]
cx = Inches(0.95)
for c in chips:
    w = Inches(0.18 + 0.092 * len(c))
    rect(s, cx, Inches(6.25), w, Inches(0.44), fill=LIGHT, line=LINE, lw=1)
    txt(s, cx, Inches(6.25), w, Inches(0.44), [[(c, 10, SLATE, HEAD, True)]],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    cx += w + Inches(0.18)

# =========================================================================
# SLIDE 2 — THE STUDY
# =========================================================================
s = slide()
header(s, "The Study", "One Factory, Four Procurement Strategies")
# left: profile
txt(s, Inches(0.78), Inches(1.7), Inches(4.4), Inches(0.3),
    [[("FACTORY A — LOAD PROFILE", 11, TEAL, HEAD, True)]])
prof = [("9.3 GWh", "Annual electricity demand", TEAL),
        ("2.43 MW", "Peak demand · flat industrial load", INK),
        ("22–110 kV", "Grid connection voltage", INK),
        ("S. Vietnam", "HCMC region · strong solar resource", AMBER)]
py = Inches(2.05)
for v, l, c in prof:
    kpi(s, Inches(0.78), py, Inches(4.35), Inches(0.92), v, c, l)
    py += Inches(1.04)
# right: 4 scenario cards
txt(s, Inches(5.55), Inches(1.7), Inches(7), Inches(0.3),
    [[("THE FOUR SCENARIOS — EACH INDEPENDENTLY SIZE-OPTIMIZED", 11, TEAL, HEAD, True)]])
cards = [
    ("CASE 4", "Solar Only", "New TOU (Decision 963). PV only — no storage. The conventional rooftop benchmark.", GREEN),
    ("CASE 1", "Solar + BESS", "Today's tariff (current TOU). Integrated PV + battery, the baseline integrated system.", TEAL),
    ("CASE 2", "Solar + BESS", "New evening-peak TOU (Decision 963, 17:30–22:30). Integrated PV + battery.", TEAL),
    ("CASE 3", "Solar + BESS", "New TOU + two-component capacity tariff (~209,459 VND/kW/mo). PV + battery.", AMBER),
]
gx0, gy0, gw, gh, gpx, gpy = Inches(5.55), Inches(2.05), Inches(3.5), Inches(1.78), Inches(0.27), Inches(0.2)
for i, (tag, title, desc, col) in enumerate(cards):
    x = gx0 + (i % 2) * (gw + gpx)
    y = gy0 + (i // 2) * (gh + gpy)
    rect(s, x, y, gw, gh, fill=CARD, line=LINE, lw=1)
    rect(s, x, y, gw, Inches(0.08), fill=col)
    txt(s, x + Inches(0.22), y + Inches(0.2), gw - Inches(0.4), Inches(0.3),
        [[(tag, 9.5, col, HEAD, True)]])
    txt(s, x + Inches(0.22), y + Inches(0.44), gw - Inches(0.4), Inches(0.4),
        [[(title, 16, INK, HEAD, True)]])
    txt(s, x + Inches(0.22), y + Inches(0.86), gw - Inches(0.42), Inches(0.85),
        [[(desc, 10, SLATE, BODY, False)]], line_spacing=1.08)
txt(s, Inches(0.78), Inches(6.55), Inches(11.8), Inches(0.4),
    [[("Method:  ", 9.5, INK, HEAD, True),
      ("PV and battery sizes are optimized in REopt for each tariff, then evaluated over a 25-year ESCO "
       "pro forma (70% debt). All monetary figures in USD. One anomalous 24-hour data spike was removed from the metered load.",
       9.5, SLATE, BODY, False)]], line_spacing=1.1)
footer(s, 2)

# =========================================================================
# SLIDE 3 — FINDING 1: SOLAR-ONLY CEILING
# =========================================================================
s = slide()
header(s, "Finding 1", "Solar Alone Stalls at One-Third of Demand", accent=GREEN)
c4 = M["case_4"]
# left narrative
txt(s, Inches(0.78), Inches(1.75), Inches(5.7), Inches(0.5),
    [[("Even fully optimized, a solar-only system meets only ", 14, SLATE, BODY, False),
      ("36% of the factory's load.", 14, INK, HEAD, True)]], line_spacing=1.25)
txt(s, Inches(0.78), Inches(2.75), Inches(5.75), Inches(3.0),
    [[("Why it plateaus", 12, GREEN, HEAD, True)],
     [("Under Decision 963 the priced peak shifts to 17:30–22:30 — after sunset. "
       "Daytime panels offset only midday consumption.", 11.5, SLATE, BODY, False)],
     [("The expensive evening peak stays on full-price grid power.", 11.5, SLATE, BODY, False)],
     [("Adding more panels doesn't help: past ~36% of load, midday surplus is "
       "curtailed or exported at low value — so the optimizer stops at 3.45 MW.", 11.5, SLATE, BODY, False)],
     [("It is a partial daytime hedge — not a path to 24/7 clean energy.", 11.5, INK, HEAD, True)]],
    line_spacing=1.18, space_after=8)
# right: doughnut + kpis
cd = CategoryChartData(); cd.categories = ["Clean (solar)", "Grid"]
cd.add_series("mix", (c4["self_suff"], round(100 - c4["self_suff"], 1)))
gf = s.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, Inches(7.0), Inches(1.75),
                        Inches(4.0), Inches(3.5), cd).chart
gf.has_title = False; gf.has_legend = True
gf.legend.position = XL_LEGEND_POSITION.BOTTOM; gf.legend.include_in_layout = False
gf.legend.font.size = Pt(11); gf.legend.font.name = BODY
plot = gf.plots[0]; plot.has_data_labels = True
dl = plot.data_labels; dl.number_format = '0.0"%"'; dl.number_format_is_linked = False
dl.font.size = Pt(12); dl.font.bold = True; dl.font.name = HEAD; dl.font.color.rgb = WHITE
ser = gf.series[0]
for i, col in enumerate([GREEN, RGBColor(0xC8, 0xD2, 0xDA)]):
    ser.points[i].format.fill.solid(); ser.points[i].format.fill.fore_color.rgb = col
try:
    gf.plots[0].doughnut_hole_size = 58
except Exception:
    pass
txt(s, Inches(7.0), Inches(5.35), Inches(4.0), Inches(0.3),
    [[("Optimal system: 3.45 MW PV, no storage", 10.5, MUTE, BODY, False)]], align=PP_ALIGN.CENTER)
# bottom KPI strip
kpis = [(f"{c4['self_suff']:.0f}%", "Clean self-supply", GREEN),
        (f"${c4['bill_sav']/1000:.0f}k", "Annual bill savings", INK),
        (f"{c4['eq_irr']:.1f}%", "Equity IRR", TEAL),
        (f"${c4['npv']/1e6:.2f}M", "25-year NPV", INK)]
kx = Inches(0.78)
for v, l, c in kpis:
    kpi(s, kx, Inches(5.95), Inches(2.78), Inches(0.92), v, c, l)
    kx += Inches(2.93)
footer(s, 3)

# =========================================================================
# SLIDE 4 — FINDING 2: BESS DOUBLES IT
# =========================================================================
s = slide()
header(s, "Finding 2", "Storage Nearly Doubles Clean Supply — and Savings", accent=TEAL)
c2 = M["case_2"]
txt(s, Inches(0.78), Inches(1.62), Inches(11.7), Inches(0.5),
    [[("Same factory, same new tariff — add a battery (Case 4 → Case 2):", 14, INK, HEAD, True)]])
# two charts
cd1 = CategoryChartData(); cd1.categories = ["Solar Only", "Solar + BESS"]
cd1.add_series("s", (c4["self_suff"], c2["self_suff"]))
ch1 = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.9), Inches(2.35),
                         Inches(5.2), Inches(3.2), cd1).chart
style_chart(ch1, [RGBColor(0xB9, 0xCF, 0xC2), TEAL], num_fmt='0.0"%"', max_scale=80)
txt(s, Inches(0.9), Inches(2.1), Inches(5.2), Inches(0.3),
    [[("CLEAN SELF-SUPPLY (% of load)", 11, TEAL, HEAD, True)]], align=PP_ALIGN.CENTER)

cd2 = CategoryChartData(); cd2.categories = ["Solar Only", "Solar + BESS"]
cd2.add_series("s", (round(c4["bill_sav"]/1000), round(c2["bill_sav"]/1000)))
ch2 = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(6.55), Inches(2.35),
                         Inches(5.2), Inches(3.2), cd2).chart
style_chart(ch2, [RGBColor(0xCF, 0xC4, 0xB0), AMBER], num_fmt='"$"0"k"', max_scale=700)
txt(s, Inches(6.55), Inches(2.1), Inches(5.2), Inches(0.3),
    [[("ANNUAL BILL SAVINGS (USD)", 11, AMBER, HEAD, True)]], align=PP_ALIGN.CENTER)
# callout strip
rect(s, Inches(0.78), Inches(5.95), Inches(11.78), Inches(0.95), fill=LIGHT, line=LINE, lw=1)
rect(s, Inches(0.78), Inches(5.95), Inches(0.09), Inches(0.95), fill=TEAL)
gain = c2["self_suff"] - c4["self_suff"]
txt(s, Inches(1.08), Inches(6.06), Inches(11.2), Inches(0.8),
    [[(f"+{gain:.0f} points clean energy   ·   +${(c2['bill_sav']-c4['bill_sav'])/1000:.0f}k/yr savings   ·   ",
       13.5, INK, HEAD, True),
      (f"and the package still earns {c2['eq_irr']:.0f}% equity IRR with ${c2['npv']/1e6:.2f}M NPV.",
       13.5, SLATE, BODY, False)],
     [("The battery shifts midday solar into the evening peak — turning a daytime hedge into round-the-clock clean supply.",
       11, SLATE, BODY, False)]], anchor=MSO_ANCHOR.MIDDLE, space_after=4, line_spacing=1.1)
footer(s, 4)

# =========================================================================
# SLIDE 5 — FINDING 3: TWO-COMPONENT TARIFF
# =========================================================================
s = slide()
header(s, "Finding 3", "The 2026 Capacity Tariff Makes Peak-Shaving the Prize", accent=AMBER)
c3 = M["case_3"]
txt(s, Inches(0.78), Inches(1.7), Inches(5.7), Inches(3.4),
    [[("When EVN starts billing capacity (~209,459 VND/kW/month), the dominant "
       "value flips from energy arbitrage to ", 12, SLATE, BODY, False),
      ("demand-charge avoidance.", 12, INK, HEAD, True)],
     [("", 6, SLATE, BODY, False)],
     [("The optimizer adds battery power to cut the grid peak nearly in half:", 12, SLATE, BODY, False)],
     [(f"{c3['peak_grid_bau']:,} kW → {c3['peak_grid']:,} kW", 20, AMBER, HEAD, True),
      (f"  (−{round((1-c3['peak_grid']/c3['peak_grid_bau'])*100)}%)", 14, SLATE, HEAD, True)],
     [("", 6, SLATE, BODY, False)],
     [("Returns tighten as the system grows — but every metric stays bankable, "
       "and this is the tariff Factory A will actually face in 2026.", 12, SLATE, BODY, False)]],
    line_spacing=1.15, space_after=6)
# chart: peak grid demand BAU vs optimized
cd = CategoryChartData(); cd.categories = ["Grid peak\n(no battery)", "Grid peak\n(with BESS)"]
cd.add_series("s", (c3["peak_grid_bau"], c3["peak_grid"]))
ch = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(6.7), Inches(2.0),
                        Inches(5.1), Inches(3.0), cd).chart
style_chart(ch, [RED, AMBER], num_fmt='#,##0" kW"', max_scale=2800)
txt(s, Inches(6.7), Inches(1.78), Inches(5.1), Inches(0.3),
    [[("PEAK GRID DEMAND — CASE 3", 11, AMBER, HEAD, True)]], align=PP_ALIGN.CENTER)
# KPI strip
kpis = [(f"${c3['dem_sav']/1000:.0f}k", "Demand-charge\nsavings / yr", AMBER),
        (f"-65%", "Total bill\n($758k → $263k)", TEAL),
        (f"{c3['eq_irr']:.1f}%", "Equity IRR\n(still bankable)", INK),
        (f"{c3['dscr']:.2f}", "Avg DSCR\n(debt-serviceable)", INK)]
kx = Inches(0.78)
for v, l, c in kpis:
    rect(s, kx, Inches(5.7), Inches(2.78), Inches(1.18), fill=CARD, line=LINE, lw=1)
    rect(s, kx, Inches(5.7), Inches(0.07), Inches(1.18), fill=c)
    txt(s, kx + Inches(0.24), Inches(5.85), Inches(2.45), Inches(0.9),
        [[(v, 23, c, HEAD, True)], [(l, 10, SLATE, BODY, False)]],
        anchor=MSO_ANCHOR.MIDDLE, space_after=2, line_spacing=1.0)
    kx += Inches(2.93)
footer(s, 5)

# =========================================================================
# SLIDE 6 — ECONOMICS DASHBOARD (TABLE)
# =========================================================================
s = slide()
header(s, "Side by Side", "Full Techno-Economic Comparison", accent=INK)
rows = [
    ("", "Solar Only", "Solar + BESS", "Solar + BESS", "Solar + BESS"),
    ("Tariff regime", "New TOU 963", "Current TOU", "New TOU 963", "TOU 963 + 2-comp."),
    ("PV size", "3.45 MW", "5.32 MW", "5.91 MW", "5.77 MW"),
    ("Battery (power / energy)", "—", "1.66 / 8.3", "1.80 / 10.7", "1.83 / 11.7"),
    ("Clean self-supply", "35.8%", "59.5%", "65.5%", "65.8%"),
    ("Total CapEx", "$1.66M", "$3.68M", "$4.27M", "$4.32M"),
    ("Annual bill savings", "$245k", "$531k", "$569k", "$494k"),
    ("Equity IRR", "18.7%", "18.2%", "16.1%", "12.4%"),
    ("25-year NPV", "$0.80M", "$1.65M", "$1.44M", "$0.59M"),
    ("Average DSCR", "1.33", "1.31", "1.21", "1.01"),
    ("Simple payback", "9.0 yr", "9.4 yr", "10.5 yr", "12.2 yr"),
]
nr, nc = len(rows), 5
tx, ty, tw, th = Inches(0.78), Inches(1.65), Inches(9.0), Inches(5.05)
gtbl = s.shapes.add_table(nr, nc, tx, ty, tw, th).table
gtbl.first_row = False; gtbl.horz_banding = False
gtbl.columns[0].width = Inches(2.7)
for ci in range(1, 5):
    gtbl.columns[ci].width = Inches(1.575)
hl = [None, GREEN, TEAL, TEAL, AMBER]
for ri, row in enumerate(rows):
    gtbl.rows[ri].height = Inches(0.43) if ri >= 2 else Inches(0.4)
    for ci, val in enumerate(row):
        cell = gtbl.cell(ri, ci)
        cell.margin_left = Inches(0.12); cell.margin_right = Inches(0.06)
        cell.margin_top = Inches(0.02); cell.margin_bottom = Inches(0.02)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf = cell.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
        r = p.add_run(); r.text = val
        # fills
        if ri == 0:
            cell.fill.solid(); cell.fill.fore_color.rgb = INK if ci > 0 else WHITE
            if ci == 0:
                cell.fill.background()
            else:
                cell.fill.fore_color.rgb = hl[ci]
            r.font.size = Pt(12); r.font.bold = True; r.font.name = HEAD
            r.font.color.rgb = WHITE
        elif ri == 1:
            cell.fill.solid(); cell.fill.fore_color.rgb = LIGHT
            r.font.size = Pt(10); r.font.name = BODY; r.font.italic = True
            r.font.color.rgb = SLATE; r.font.bold = (ci == 0)
        else:
            cell.fill.solid(); cell.fill.fore_color.rgb = WHITE if ri % 2 else CARD
            r.font.size = Pt(11.5); r.font.name = HEAD if ci == 0 else BODY
            r.font.color.rgb = INK if ci == 0 else SLATE
            r.font.bold = (ci == 0)
            # emphasize key rows
            if rows[ri][0] in ("Clean self-supply", "Equity IRR"):
                r.font.bold = True
                if ci > 0:
                    r.font.color.rgb = hl[ci]
# right rail takeaways
rx = Inches(10.05)
rect(s, rx, Inches(1.65), Inches(2.5), Inches(5.05), fill=LIGHT, line=LINE, lw=1)
txt(s, rx + Inches(0.22), Inches(1.85), Inches(2.1), Inches(0.3),
    [[("READING THE TABLE", 10.5, INK, HEAD, True)]])
notes = [("All four are NPV-positive and clear DSCR ≥ 1.0 — every option is financeable.", INK),
         ("Adding BESS lifts clean supply from 36% to ~60–66%.", TEAL),
         ("Tariff reform compresses returns (18% → 12%) but never breaks bankability.", AMBER),
         ("Solar-only has the highest IRR but the lowest absolute value and clean share.", SLATE)]
ny = Inches(2.3)
for t, c in notes:
    rect(s, rx + Inches(0.22), ny + Inches(0.04), Inches(0.07), Inches(0.07), fill=c, shape=MSO_SHAPE.OVAL)
    txt(s, rx + Inches(0.4), ny - Inches(0.04), Inches(1.95), Inches(1.1),
        [[(t, 10, SLATE, BODY, False)]], line_spacing=1.08)
    ny += Inches(1.12)
footer(s, 6)

# =========================================================================
# SLIDE 7 — CONCLUSIONS
# =========================================================================
s = slide()
header(s, "Conclusions", "What the Four Cases Tell CEBA Members", accent=TEAL)
concl = [
    ("01", "Solar alone is only a partial hedge", GREEN,
     "Optimized rooftop PV caps at ~36% of demand and leaves the priced evening peak on the grid."),
    ("02", "BESS is what unlocks 24/7 clean energy", TEAL,
     "Storage lifts clean self-supply to 60–66% and roughly doubles bill savings — while equity IRR stays 16–18% and NPV stays positive."),
    ("03", "Tariff reform compresses but never breaks the case", INK,
     "Across current TOU, Decision 963 and the two-component tariff, every Solar+BESS case clears DSCR ≥ 1.0 and remains financeable."),
    ("04", "Size the battery for the 2026 tariff — not today's", AMBER,
     "Under the two-component capacity charge, demand-charge avoidance (peak −46%) becomes the single largest value stream. Size for power, not just energy."),
]
cy = Inches(1.72)
for tag, title, col, desc in concl:
    rect(s, Inches(0.78), cy, Inches(8.05), Inches(1.06), fill=CARD, line=LINE, lw=1)
    rect(s, Inches(0.78), cy, Inches(0.78), Inches(1.06), fill=col)
    txt(s, Inches(0.78), cy + Inches(0.1), Inches(0.78), Inches(0.86),
        [[(tag, 22, WHITE, HEAD, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(1.78), cy + Inches(0.14), Inches(6.9), Inches(0.4),
        [[(title, 14.5, INK, HEAD, True)]])
    txt(s, Inches(1.78), cy + Inches(0.55), Inches(6.95), Inches(0.5),
        [[(desc, 10.5, SLATE, BODY, False)]], line_spacing=1.08)
    cy += Inches(1.15)
# recommendation rail
rx = Inches(9.05)
rect(s, rx, Inches(1.72), Inches(3.5), Inches(4.74), fill=INK)
rect(s, rx, Inches(1.72), Inches(3.5), Inches(0.1), fill=AMBER)
txt(s, rx + Inches(0.3), Inches(2.05), Inches(2.95), Inches(0.4),
    [[("RECOMMENDATION", 11, AMBER, HEAD, True)]])
txt(s, rx + Inches(0.3), Inches(2.5), Inches(2.95), Inches(2.6),
    [[("Pursue integrated ", 15, WHITE, BODY, False),
      ("Solar + BESS now", 15, WHITE, HEAD, True),
      (" — and size storage to the ", 15, WHITE, BODY, False),
      ("tariff you will face in 2026", 15, AMBER, HEAD, True),
      (", not the one you have today.", 15, WHITE, BODY, False)]], line_spacing=1.25)
txt(s, rx + Inches(0.3), Inches(5.35), Inches(2.95), Inches(1.0),
    [[("The battery converts a daytime solar hedge into bankable, round-the-clock clean supply.",
       11, RGBColor(0xC5, 0xD6, 0xDC), BODY, False)]], line_spacing=1.18)
txt(s, Inches(0.78), Inches(6.62), Inches(11.8), Inches(0.4),
    [[("Source: REopt system optimization + Vietnam ESCO pro forma (25 yr, 70% debt). Figures USD. "
       "One anomalous 24-hour load spike removed prior to modeling.", 8.5, MUTE, BODY, False)]])
footer(s, 7)

out = "Factory_A_Solar_BESS_Case_Study.pptx"
prs.save(out)
print("saved", out)
